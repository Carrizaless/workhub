"""Generate the WorkHub Invite PPTX presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BLUE = RGBColor(0x25, 0x63, 0xEB)       # --accent
DARK = RGBColor(0x11, 0x18, 0x27)        # --foreground
GRAY = RGBColor(0x6B, 0x72, 0x80)        # --muted
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)   # --muted-bg
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)  # --accent-light

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_shape_with_text(slide, left, top, width, height, text, bg_color, text_color, font_size=14, bold=False, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg_color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = 'Calibri'
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    # vertically center
    tf.auto_size = None
    from pptx.enum.text import MSO_ANCHOR
    tf.word_wrap = True
    shp.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shp

# ─── SLIDE 1: Cover ───
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)

# Blue accent bar at top
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_text_box(slide, 1, 1.5, 11.3, 1, 'WorkHub', font_size=52, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.6, 11.3, 1, 'Plataforma de Gestion de Trabajo Academico', font_size=28, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

# Divider line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = BLUE
line.line.fill.background()

add_text_box(slide, 1, 4, 11.3, 1.2, 'Unete como colaborador y genera ingresos\ncompletando tareas academicas.', font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# CTA badge
add_shape_with_text(slide, 4.8, 5.5, 3.7, 0.7, 'Comienza Ahora', BLUE, WHITE, font_size=20, bold=True)

# ─── SLIDE 2: Que es WorkHub? ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_text_box(slide, 0.8, 0.5, 11.7, 0.8, 'Que es WorkHub?', font_size=36, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 1.4, 10.3, 1.2,
    'WorkHub es una plataforma donde administradores publican tareas academicas\n'
    'y colaboradores las aceptan, completan y reciben pago por su trabajo.',
    font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER)

# Feature cards
features = [
    ('Tareas Variadas', 'Ensayos, investigaciones,\npresentaciones, y mas.'),
    ('Pago por Tarea', 'Cobra por cada tarea\ncompletada y aprobada.'),
    ('Chat Integrado', 'Comunicate directamente\ncon el administrador.'),
    ('Dashboard Personal', 'Monitorea tu progreso,\ncalificaciones y ganancias.'),
]

card_w = 2.6
gap = 0.3
total_w = len(features) * card_w + (len(features) - 1) * gap
start_x = (13.333 - total_w) / 2

for i, (title, desc) in enumerate(features):
    x = start_x + i * (card_w + gap)
    y = 3.0
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(3.2))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.fill.background()

    # Icon circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + card_w/2 - 0.35), Inches(y + 0.4), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = LIGHT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = BLUE
    p.font.bold = True

    add_text_box(slide, x + 0.2, y + 1.3, card_w - 0.4, 0.5, title, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, y + 1.8, card_w - 0.4, 1.2, desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 3: Beneficios ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_text_box(slide, 0.8, 0.5, 11.7, 0.8, 'Beneficios para Colaboradores', font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

benefits = [
    ('Flexibilidad Total', 'Trabaja desde cualquier lugar,\nen tu propio horario.'),
    ('Ingresos Constantes', 'Recibe pago por cada tarea\naprobada directamente en tu billetera.'),
    ('Desarrollo Profesional', 'Mejora tus habilidades\nacademicas y de investigacion.'),
    ('Soporte Directo', 'Chat integrado para resolver\ndudas al instante con el admin.'),
    ('Gestion de Archivos', 'Sube y descarga documentos\nde forma segura en la plataforma.'),
    ('Sin Compromisos', 'Acepta solo las tareas\nque te interesen.'),
]

cols = 3
rows = 2
card_w = 3.4
card_h = 2.0
gap_x = 0.4
gap_y = 0.4
total_w = cols * card_w + (cols - 1) * gap_x
start_x = (13.333 - total_w) / 2
start_y = 1.8

for idx, (title, desc) in enumerate(benefits):
    col = idx % cols
    row = idx // cols
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(card_h))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    card.line.color.rgb = RGBColor(0x26, 0x26, 0x26)
    card.line.width = Pt(1)

    add_text_box(slide, x + 0.3, y + 0.3, card_w - 0.6, 0.5, title, font_size=17, color=BLUE, bold=True, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, x + 0.3, y + 0.9, card_w - 0.6, 1.0, desc, font_size=13, color=GRAY, alignment=PP_ALIGN.LEFT)


# ─── SLIDE 4: Como unirte ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_text_box(slide, 0.8, 0.5, 11.7, 0.8, 'Como Unirte?', font_size=36, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

join_steps = [
    ('1', 'Recibe el Enlace', 'El administrador te comparte\nel enlace de acceso a WorkHub.'),
    ('2', 'Registrate', 'Crea tu cuenta con tu email\ny contrasena en segundos.'),
    ('3', 'Explora Tareas', 'Revisa las tareas disponibles\ny acepta las que te interesen.'),
    ('4', 'Cobra tu Trabajo', 'Completa tareas, recibe\naprobacion y cobra.'),
]

card_w = 2.7
gap = 0.3
total_w = len(join_steps) * card_w + (len(join_steps) - 1) * gap
start_x = (13.333 - total_w) / 2

for i, (num, title, desc) in enumerate(join_steps):
    x = start_x + i * (card_w + gap)
    y = 2.0

    # Card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(3.5))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.fill.background()

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + card_w/2 - 0.4), Inches(y + 0.4), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True

    add_text_box(slide, x + 0.2, y + 1.5, card_w - 0.4, 0.5, title, font_size=17, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, y + 2.1, card_w - 0.4, 1.2, desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < len(join_steps) - 1:
        arrow_x = x + card_w + 0.02
        add_text_box(slide, arrow_x, y + 1.5, gap - 0.04, 0.5, '>', font_size=24, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 5: CTA / Cierre ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_text_box(slide, 1, 2, 11.3, 1, 'Listo para empezar?', font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 2, 3.2, 9.3, 1.2,
    'Unete a WorkHub como colaborador y comienza a generar\ningresos completando tareas academicas.',
    font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

add_shape_with_text(slide, 4.5, 4.8, 4.3, 0.8, 'Solicita tu Invitacion', BLUE, WHITE, font_size=22, bold=True)

add_text_box(slide, 1, 6.2, 11.3, 0.5, 'WorkHub  |  Gestion de Trabajo Academico', font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# Save
out_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
out_path = os.path.join(out_dir, 'WorkHub_Invitacion_Colaboradores.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
