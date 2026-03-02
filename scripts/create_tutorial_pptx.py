"""Generate the WorkHub Tutorial PPTX presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

BLUE = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x11, 0x18, 0x27)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF3, 0xF4, 0xF6)
LIGHT_BLUE = RGBColor(0xDB, 0xEA, 0xFE)
GREEN = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=14, color=DARK, alignment=PP_ALIGN.LEFT):
    """Add a text box with multiple paragraphs."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, sz, clr, bld) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = 'Calibri'
        p.alignment = alignment
        p.space_after = Pt(4)
    return txBox

def add_shape_with_text(slide, left, top, width, height, text, bg_color, text_color, font_size=14, bold=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = bg_color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER
    return shp


# ─── SLIDE 1: Cover ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_text_box(slide, 1, 1.5, 11.3, 1, 'WorkHub', font_size=52, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.6, 11.3, 1, 'Tutorial: Flujo de Trabajo', font_size=28, color=WHITE, bold=False, alignment=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.6), Inches(2.3), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = BLUE
line.line.fill.background()

add_text_box(slide, 1, 4.0, 11.3, 1.2,
    'Guia paso a paso del ciclo completo de una tarea:\ndesde la publicacion hasta el pago.',
    font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)

# Role legend
add_shape_with_text(slide, 3.5, 5.5, 1.8, 0.5, 'Admin', BLUE, WHITE, font_size=14, bold=True)
add_shape_with_text(slide, 5.7, 5.5, 1.8, 0.5, 'Colaborador', GREEN, WHITE, font_size=14, bold=True)
add_shape_with_text(slide, 7.9, 5.5, 1.8, 0.5, 'Sistema', AMBER, WHITE, font_size=14, bold=True)


# ─── SLIDE 2: Paso 1 - Admin publica tarea ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_shape_with_text(slide, 0.8, 0.4, 1.2, 0.6, 'Paso 1', BLUE, WHITE, font_size=18, bold=True)
add_text_box(slide, 2.3, 0.4, 8, 0.7, 'El Admin Publica una Tarea', font_size=32, color=DARK, bold=True)
add_shape_with_text(slide, 10.8, 0.45, 1.3, 0.5, 'Admin', BLUE, WHITE, font_size=13, bold=True)

add_text_box(slide, 0.8, 1.4, 11.5, 0.8,
    'El administrador crea una nueva tarea desde el panel. Define todos los detalles\npara que los colaboradores puedan evaluarla y decidir si la aceptan.',
    font_size=16, color=GRAY, alignment=PP_ALIGN.LEFT)

# Detail cards
details = [
    ('Titulo y Descripcion', 'Define claramente que se necesita\nhacer en la tarea.'),
    ('Categoria', 'Clasifica la tarea: ensayo,\ninvestigacion, presentacion, etc.'),
    ('Fecha Limite', 'Establece la fecha maxima\nde entrega para la tarea.'),
    ('Monto de Pago', 'Define cuanto se pagara al\ncolaborador por completarla.'),
]

card_w = 2.8
gap = 0.35
total_w = len(details) * card_w + (len(details) - 1) * gap
start_x = (13.333 - total_w) / 2

for i, (title, desc) in enumerate(details):
    x = start_x + i * (card_w + gap)
    y = 2.8

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(2.8))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.fill.background()

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + card_w/2 - 0.35), Inches(y + 0.4), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = LIGHT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = BLUE
    p.font.bold = True

    add_text_box(slide, x + 0.2, y + 1.3, card_w - 0.4, 0.5, title, font_size=15, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, y + 1.9, card_w - 0.4, 1.0, desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 6.2, 11.5, 0.5,
    'La tarea queda publicada y visible para todos los colaboradores.',
    font_size=14, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 3: Paso 2 - Colaborador acepta ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_shape_with_text(slide, 0.8, 0.4, 1.2, 0.6, 'Paso 2', GREEN, WHITE, font_size=18, bold=True)
add_text_box(slide, 2.3, 0.4, 8, 0.7, 'El Colaborador Acepta la Tarea', font_size=32, color=WHITE, bold=True)
add_shape_with_text(slide, 10.5, 0.45, 1.8, 0.5, 'Colaborador', GREEN, WHITE, font_size=13, bold=True)

add_text_box(slide, 0.8, 1.4, 11.5, 0.8,
    'El colaborador revisa las tareas disponibles, lee los detalles\ny decide aceptar la que le interese. Nadie le obliga a tomar ninguna.',
    font_size=16, color=GRAY, alignment=PP_ALIGN.LEFT)

# Process flow
flow_steps = [
    ('Revisa Tareas', 'Explora las tareas\ndisponibles en el panel.'),
    ('Lee los Detalles', 'Revisa descripcion, fecha\nlimite y monto de pago.'),
    ('Acepta la Tarea', 'Confirma que desea\ntrabajar en ella.'),
    ('Comienza a Trabajar', 'La tarea pasa a su\nlista de pendientes.'),
]

card_w = 2.6
gap = 0.35
total_w = len(flow_steps) * card_w + (len(flow_steps) - 1) * gap
start_x = (13.333 - total_w) / 2

for i, (title, desc) in enumerate(flow_steps):
    x = start_x + i * (card_w + gap)
    y = 2.8

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(2.8))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    card.line.color.rgb = RGBColor(0x26, 0x26, 0x26)
    card.line.width = Pt(1)

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + card_w/2 - 0.35), Inches(y + 0.35), Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = GREEN
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True

    add_text_box(slide, x + 0.2, y + 1.3, card_w - 0.4, 0.5, title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, y + 1.9, card_w - 0.4, 1.0, desc, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(flow_steps) - 1:
        arrow_x = x + card_w + 0.02
        add_text_box(slide, arrow_x, y + 1.2, gap - 0.04, 0.5, '>', font_size=24, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 6.2, 11.5, 0.5,
    'El colaborador tiene total libertad de aceptar o rechazar tareas.',
    font_size=14, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 4: Paso 3 - Trabajo y entrega ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_shape_with_text(slide, 0.8, 0.4, 1.2, 0.6, 'Paso 3', GREEN, WHITE, font_size=18, bold=True)
add_text_box(slide, 2.3, 0.4, 8, 0.7, 'Trabajo, Chat y Entrega', font_size=32, color=DARK, bold=True)
add_shape_with_text(slide, 10.5, 0.45, 1.8, 0.5, 'Colaborador', GREEN, WHITE, font_size=13, bold=True)

add_text_box(slide, 0.8, 1.4, 11.5, 0.8,
    'Durante el trabajo, el colaborador puede comunicarse con el admin via chat\ny subir archivos. Al terminar, entrega su trabajo desde la plataforma.',
    font_size=16, color=GRAY, alignment=PP_ALIGN.LEFT)

# Three columns
cols_data = [
    ('Comunicacion', [
        'Chat integrado por tarea',
        'Preguntas y aclaraciones',
        'Notificaciones en tiempo real',
        'Adjuntar archivos al chat',
    ]),
    ('Trabajo', [
        'Acceso a los detalles de la tarea',
        'Subir archivos de trabajo',
        'Soporte de PDF, Word, imagenes',
        'Gestionar multiples tareas',
    ]),
    ('Entrega', [
        'Boton de entregar tarea',
        'Adjuntar archivos finales',
        'Comentarios de entrega',
        'Confirmacion automatica',
    ]),
]

card_w = 3.6
gap = 0.35
total_w = len(cols_data) * card_w + (len(cols_data) - 1) * gap
start_x = (13.333 - total_w) / 2

for i, (title, items) in enumerate(cols_data):
    x = start_x + i * (card_w + gap)
    y = 2.8

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(3.6))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.fill.background()

    add_text_box(slide, x + 0.3, y + 0.3, card_w - 0.6, 0.5, title, font_size=18, color=DARK, bold=True, alignment=PP_ALIGN.LEFT)

    # Divider
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.3), Inches(y + 0.9), Inches(card_w - 0.6), Inches(0.03))
    div.fill.solid()
    div.fill.fore_color.rgb = BLUE
    div.line.fill.background()

    for j, item in enumerate(items):
        item_y = y + 1.1 + j * 0.55
        # Bullet dot
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.4), Inches(item_y + 0.07), Inches(0.12), Inches(0.12))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        add_text_box(slide, x + 0.65, item_y, card_w - 1.0, 0.4, item, font_size=13, color=GRAY, alignment=PP_ALIGN.LEFT)


# ─── SLIDE 5: Paso 4 - Revision y aprobacion ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_shape_with_text(slide, 0.8, 0.4, 1.2, 0.6, 'Paso 4', BLUE, WHITE, font_size=18, bold=True)
add_text_box(slide, 2.3, 0.4, 8, 0.7, 'Revision y Aprobacion', font_size=32, color=WHITE, bold=True)
add_shape_with_text(slide, 10.8, 0.45, 1.3, 0.5, 'Admin', BLUE, WHITE, font_size=13, bold=True)

add_text_box(slide, 0.8, 1.4, 11.5, 0.8,
    'El administrador revisa la entrega del colaborador. Tiene tres opciones\nsobre que hacer con la tarea entregada.',
    font_size=16, color=GRAY, alignment=PP_ALIGN.LEFT)

# Three outcome cards
outcomes = [
    ('Aprobar', 'La tarea cumple con los\nrequisitos. Se procede\nal pago automatico.', GREEN, RGBColor(0x06, 0x5F, 0x46)),
    ('Solicitar Correccion', 'La tarea necesita ajustes.\nEl colaborador recibe\ncomentarios y corrige.', AMBER, RGBColor(0x92, 0x40, 0x0E)),
    ('Rechazar', 'La tarea no cumple.\nSe notifica al colaborador\ncon el motivo.', RGBColor(0xEF, 0x44, 0x44), RGBColor(0x99, 0x1B, 0x1B)),
]

card_w = 3.4
gap = 0.5
total_w = len(outcomes) * card_w + (len(outcomes) - 1) * gap
start_x = (13.333 - total_w) / 2

for i, (title, desc, accent, dark_accent) in enumerate(outcomes):
    x = start_x + i * (card_w + gap)
    y = 2.8

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(3.2))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    card.line.color.rgb = accent
    card.line.width = Pt(2)

    # Top accent bar
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x + 0.15), Inches(y + 0.15), Inches(card_w - 0.3), Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent
    top_bar.line.fill.background()

    add_text_box(slide, x + 0.3, y + 0.5, card_w - 0.6, 0.5, title, font_size=20, color=accent, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.3, y + 1.3, card_w - 0.6, 1.5, desc, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 6: Paso 5 - Pago ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_shape_with_text(slide, 0.8, 0.4, 1.2, 0.6, 'Paso 5', AMBER, WHITE, font_size=18, bold=True)
add_text_box(slide, 2.3, 0.4, 8, 0.7, 'Pago Automatico', font_size=32, color=DARK, bold=True)
add_shape_with_text(slide, 10.6, 0.45, 1.5, 0.5, 'Sistema', AMBER, WHITE, font_size=13, bold=True)

add_text_box(slide, 0.8, 1.4, 11.5, 0.8,
    'Una vez aprobada la tarea, el sistema acredita automaticamente el pago\nen la billetera del colaborador.',
    font_size=16, color=GRAY, alignment=PP_ALIGN.LEFT)

# Payment flow
pay_steps = [
    ('Tarea Aprobada', 'El admin aprueba la entrega.'),
    ('Pago Acreditado', 'El monto se suma a la\nbilletera del colaborador.'),
    ('Disponible', 'El colaborador puede\nver su saldo actualizado.'),
]

card_w = 3.2
gap = 0.6
total_w = len(pay_steps) * card_w + (len(pay_steps) - 1) * gap
start_x = (13.333 - total_w) / 2

for i, (title, desc) in enumerate(pay_steps):
    x = start_x + i * (card_w + gap)
    y = 3.0

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(2.5))
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_BG
    card.line.fill.background()

    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + card_w/2 - 0.4), Inches(y + 0.3), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = AMBER
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True

    add_text_box(slide, x + 0.2, y + 1.3, card_w - 0.4, 0.5, title, font_size=16, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + 0.2, y + 1.8, card_w - 0.4, 0.8, desc, font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(pay_steps) - 1:
        arrow_x = x + card_w + 0.05
        add_text_box(slide, arrow_x, y + 1.0, gap - 0.1, 0.5, '>', font_size=28, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 7: Resumen del flujo ───
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE
bar.line.fill.background()

add_text_box(slide, 0.8, 0.4, 11.7, 0.8, 'Resumen del Flujo Completo', font_size=34, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Horizontal flow
flow = [
    ('Admin Publica', BLUE),
    ('Colaborador\nAcepta', GREEN),
    ('Trabajo y\nEntrega', GREEN),
    ('Admin\nRevisa', BLUE),
    ('Pago\nAutomatico', AMBER),
]

box_w = 2.0
gap = 0.35
total_w = len(flow) * box_w + (len(flow) - 1) * gap
start_x = (13.333 - total_w) / 2
y = 2.2

for i, (label, color) in enumerate(flow):
    x = start_x + i * (box_w + gap)

    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(box_w), Inches(1.3))
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Calibri'
    p.alignment = PP_ALIGN.CENTER

    if i < len(flow) - 1:
        arrow_x = x + box_w + 0.02
        add_text_box(slide, arrow_x, y + 0.35, gap - 0.04, 0.5, '>', font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Key points
add_text_box(slide, 0.8, 4.2, 11.7, 0.6, 'Puntos Clave', font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

points = [
    'El admin publica tareas, no las asigna. El colaborador elige cuales aceptar.',
    'El chat integrado permite comunicacion fluida durante todo el proceso.',
    'El pago se acredita automaticamente en la billetera al aprobar la tarea.',
    'El colaborador tiene total libertad y transparencia sobre su trabajo.',
]

for i, point in enumerate(points):
    py = 5.0 + i * 0.5
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3.0), Inches(py + 0.07), Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = BLUE
    dot.line.fill.background()
    add_text_box(slide, 3.3, py, 7.0, 0.4, point, font_size=14, color=GRAY, alignment=PP_ALIGN.LEFT)

add_text_box(slide, 1, 7.0, 11.3, 0.4, 'WorkHub  |  Gestion de Trabajo Academico', font_size=13, color=GRAY, alignment=PP_ALIGN.CENTER)


# Save
out_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
out_path = os.path.join(out_dir, 'WorkHub_Tutorial_Flujo_de_Trabajo.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
